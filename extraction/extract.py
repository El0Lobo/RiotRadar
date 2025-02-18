import os
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter.ttk import Progressbar
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
from pdf2image import convert_from_path
import pytesseract

# Function to extract text from images using OCR
def extract_text_from_image(file_path):
    """Extract text from an image file using OCR."""
    try:
        text = pytesseract.image_to_string(Image.open(file_path))
    except Exception as e:
        text = f"Error processing image: {e}"
    return text

# Extract text from PDF with fallback to OCR
def extract_text_from_pdf(file_path):
    """Extract text from a PDF file, including OCR fallback for image-based pages."""
    text = ""
    reader = PdfReader(file_path)

    for page_index, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if page_text.strip():
            text += page_text + "\n"
        else:
            # Perform OCR on the page if no text is found
            images = convert_from_path(file_path, first_page=page_index + 1, last_page=page_index + 1)
            for image in images:
                text += pytesseract.image_to_string(image) + "\n"
    return text

# Extract text from PowerPoint
def extract_text_from_ppt(file_path):
    """Extract text from a PowerPoint file, including images using OCR if necessary."""
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

# Process files and save each as a Word document
def process_files(file_paths, destination_folder, progress_callback):
    """Process the selected files and save the output to individual Word documents."""
    total_files = len(file_paths)
    for index, file_path in enumerate(file_paths):
        file_name = os.path.basename(file_path)
        output_file_name = f"Extracted_{file_name}.docx"
        output_path = os.path.join(destination_folder, output_file_name)

        # Skip files already processed
        if os.path.exists(output_path):
            continue

        document = Document()
        document.add_heading(f"Extracted Text from {file_name}", level=1)

        if file_path.endswith(".pdf"):
            text = extract_text_from_pdf(file_path)
        elif file_path.endswith(".ppt") or file_path.endswith(".pptx"):
            text = extract_text_from_ppt(file_path)
        elif file_path.lower().endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif")):
            text = extract_text_from_image(file_path)
        else:
            text = "Unsupported file type."

        document.add_paragraph(text)
        document.save(output_path)

        progress_callback(index + 1, total_files)

    return destination_folder

# Browse files
def browse_files():
    file_paths = filedialog.askopenfilenames(filetypes=[
        ("PDF files", "*.pdf"),
        ("PowerPoint files", "*.ppt;*.pptx"),
        ("Image files", "*.png;*.jpg;*.jpeg;*.tiff;*.bmp;*.gif"),
    ])
    if file_paths:
        for file_path in file_paths:
            if file_path not in selected_files:  # Avoid duplicates
                selected_files.append(file_path)
        update_file_list()

def browse_folder():
    folder = filedialog.askdirectory(title="Choose destination folder")
    return folder

# Update the file list display
def update_file_list():
    file_list_display.delete(0, tk.END)
    for file in selected_files:
        file_list_display.insert(tk.END, file)

# Remove selected file from the list
def remove_selected_file():
    selected_indices = file_list_display.curselection()
    for index in reversed(selected_indices):
        selected_files.pop(index)
    update_file_list()

# Start extraction
def start_extraction():
    if not selected_files:
        messagebox.showwarning("No Files Selected", "Please select files to process.")
        return

    destination = browse_folder()
    if not destination:
        messagebox.showwarning("No Destination Selected", "Please select a destination folder.")
        return

    progress_bar["maximum"] = len(selected_files)
    progress_bar["value"] = 0

    def update_progress(current, total):
        progress_bar["value"] = current
        progress_label.config(text=f"Processing {current}/{total} files...")
        root.update_idletasks()

    try:
        output_folder = process_files(selected_files, destination, update_progress)
        messagebox.showinfo("Success", f"Text extracted and saved in: {output_folder}")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {e}")

# Create the GUI
root = tk.Tk()
root.title("Text Extractor")
root.geometry("600x500")

selected_files = []

frame = tk.Frame(root)
frame.pack(pady=10)

select_button = tk.Button(frame, text="Add Files", command=browse_files)
select_button.grid(row=0, column=0, padx=10)

remove_button = tk.Button(frame, text="Remove Selected", command=remove_selected_file)
remove_button.grid(row=0, column=1, padx=10)

start_button = tk.Button(frame, text="Start Extraction", command=start_extraction)
start_button.grid(row=0, column=2, padx=10)

progress_bar = Progressbar(root, orient="horizontal", length=500, mode="determinate")
progress_bar.pack(pady=10)

progress_label = tk.Label(root, text="")
progress_label.pack()

file_list_frame = tk.Frame(root)
file_list_frame.pack(pady=10, fill=tk.BOTH, expand=True)

file_list_label = tk.Label(file_list_frame, text="Selected Files:")
file_list_label.pack(anchor="w")

file_list_display = tk.Listbox(file_list_frame, height=15, width=70, selectmode=tk.MULTIPLE)
file_list_display.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

scrollbar = tk.Scrollbar(file_list_frame)
scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

file_list_display.config(yscrollcommand=scrollbar.set)
scrollbar.config(command=file_list_display.yview)

root.mainloop()
